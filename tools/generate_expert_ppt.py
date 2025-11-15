from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helpers
def set_slide_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (18, 52, 86))  # dark blue
    left = Inches(0.5)
    top = Inches(0.6)
    width = Inches(12.3)
    height = Inches(1.6)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    tx = slide.shapes.add_textbox(left, top+Inches(1.2), width, Inches(0.9)).text_frame
    q = tx.paragraphs[0]
    q.text = subtitle
    q.font.size = Pt(18)
    q.font.color.rgb = RGBColor(220, 230, 240)

# Title
add_title_slide("Task Admin — Project Overview", "Features, architecture, flows and learnings")

# Slide: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (245, 247, 250))
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.6)).text_frame
title.paragraphs[0].text = "Agenda"
title.paragraphs[0].font.size = Pt(28)

body = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(5.5), Inches(4.5)).text_frame
items = ["Overview & Goals", "Architecture & Data Flow", "Key Features", "User Flows (Auth, Task CRUD, Comments)", "UX Improvements & Learnings", "Next Steps"]
for it in items:
    p = body.add_paragraph()
    p.text = f"• {it}"
    p.font.size = Pt(18)

# Slide: Architecture diagram (simple boxes)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (255, 255, 255))
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5)).text_frame
title.paragraphs[0].text = "Architecture & Data Flow"
title.paragraphs[0].font.size = Pt(26)

# Add shapes: Mobile App, Back4App, Parse DB
shapes = slide.shapes
app_box = shapes.add_shape(1, Inches(1.0), Inches(1.2), Inches(3.2), Inches(1.6))
app_box.text = "Flutter App (UI)\nScreens: Login, Home, Task Detail, Add/Edit"
app_box.text_frame.paragraphs[0].font.size = Pt(12)
app_box.fill.solid()
app_box.fill.fore_color.rgb = RGBColor(230, 245, 255)

server_box = shapes.add_shape(1, Inches(4.9), Inches(1.2), Inches(3.6), Inches(1.6))
server_box.text = "Back4App / Parse Server\nREST API & Cloud Functions"
server_box.text_frame.paragraphs[0].font.size = Pt(12)
server_box.fill.solid()
server_box.fill.fore_color.rgb = RGBColor(255, 245, 230)

db_box = shapes.add_shape(1, Inches(8.5), Inches(1.2), Inches(3.2), Inches(1.6))
db_box.text = "Parse Database\n(Cloud data storage)"
db_box.text_frame.paragraphs[0].font.size = Pt(12)
db_box.fill.solid()
db_box.fill.fore_color.rgb = RGBColor(230, 255, 230)

# arrows
conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, app_box.left+app_box.width, app_box.top+app_box.height/2, server_box.left, server_box.top+server_box.height/2)
conn.line.color.rgb = RGBColor(90, 120, 200)
conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, server_box.left+server_box.width, server_box.top+server_box.height/2, db_box.left, db_box.top+db_box.height/2)
conn.line.color.rgb = RGBColor(200, 130, 60)

# Slide: Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (250, 250, 253))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "Key Features"
body = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(6.5), Inches(5.5)).text_frame
features = [
    "Authentication: Register/Login/Logout (Parse)",
    "Task CRUD: Create / Read / Update / Delete",
    "Threaded comments per task",
    "Task fields: status, start/end dates, discussion",
    "Search, filters, and pull-to-refresh",
    "Theming & consistent input styling"
]
for f in features:
    p = body.add_paragraph()
    p.text = f"• {f}"
    p.font.size = Pt(16)

# Slide: User Flow - Auth
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (255, 255, 255))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "User Flow — Authentication"
flow = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(3.5)).text_frame
steps = ["Launch → Splash (auto-login check)", "Register (email, password) → Parse User created", "Login → Auth token saved", "Home (task list)" ]
for s in steps:
    p = flow.add_paragraph()
    p.text = f"→ {s}"
    p.font.size = Pt(16)

# Slide: User Flow - Task CRUD
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (245, 247, 250))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "User Flow — Task CRUD"
flow = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(3.5)).text_frame
steps = ["Home → Add Task (title, desc, status, dates)", "Create → Task saved to Parse & initial comment created if discussion provided", "Tap Task → Open Task Detail (threaded comments)", "Edit / Delete / Toggle Complete → Server syncs"]
for s in steps:
    p = flow.add_paragraph()
    p.text = f"→ {s}"
    p.font.size = Pt(16)

# Slide: UX Improvements & Learnings
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (255, 255, 255))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "UX Improvements & Learnings"
body = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(4.0)).text_frame
points = [
    "Centralized input styling for consistent forms",
    "Fixed keyboard overlap using AnimatedPadding and viewInsets",
    "Added status/dates & validation (endDate >= startDate)",
    "Threaded comments to improve task collaboration",
    "Prefer const and analyzer-driven fixes for cleaner code",
]
for pnt in points:
    p = body.add_paragraph()
    p.text = f"• {pnt}"
    p.font.size = Pt(16)

# Slide: Screenshots / How to demo (placeholders)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (245, 247, 250))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "Demo Screenshots (replace placeholders)"
note = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(1.0)).text_frame
note.paragraphs[0].text = "Add your actual screenshots into the /screenshots folder and re-run the generator to embed them."
note.paragraphs[0].font.size = Pt(14)

# Placeholder boxes
slide.shapes.add_shape(1, Inches(0.6), Inches(2.2), Inches(4.0), Inches(2.4)).text = "[Splash Screen]"
slide.shapes.add_shape(1, Inches(5.0), Inches(2.2), Inches(4.0), Inches(2.4)).text = "[Login/Register]"
slide.shapes.add_shape(1, Inches(9.4), Inches(2.2), Inches(2.5), Inches(2.4)).text = "[Home]"

# Slide: Next Steps & Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (255, 255, 255))
slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.paragraphs[0].text = "Next Steps & Recommendations"
body = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(11.0), Inches(4.0)).text_frame
recs = [
    "Add screenshot assets and re-run generator to embed visuals",
    "Record a 2-minute demo video following the demo script in PROJECT_GUIDE.md",
    "Add automated E2E tests for core flows",
    "Consider adding push notifications and offline sync"
]
for r in recs:
    p = body.add_paragraph()
    p.text = f"• {r}"
    p.font.size = Pt(16)

# Final slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, (18, 52, 86))
slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.0)).text_frame.paragraphs[0].text = "Thank you!"
slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12.3), Inches(0.8)).text_frame.paragraphs[0].text = "Questions?"

# Save
prs.save('project_overview_expert.pptx')
print('Wrote project_overview_expert.pptx')
