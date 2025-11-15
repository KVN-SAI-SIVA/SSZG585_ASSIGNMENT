from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Helper to add a slide with title and bullet points + speaker notes

def add_bulleted_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
            p.level = 0
        else:
            p = body.add_paragraph()
            p.text = b
            p.level = 0
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Task Admin App — Project Overview"
slide.placeholders[1].text = "Project flow, features, and learnings\nPrepared for code handoff"

# Slides
add_bulleted_slide(
    "Project Summary",
    [
        "Flutter (Material 3) app with Parse Server (Back4App) backend",
        "Core features: authentication, task CRUD, threaded comments, status/dates",
        "Goal: improve UX, add Azure-DevOps-like task fields and discussion threads",
    ],
    notes="High-level description of the app and primary goals. Good to state target platforms (web first).",
)

add_bulleted_slide(
    "High-level Architecture",
    [
        "Frontend: Flutter (lib/) — screens, widgets, theme",
        "Backend: Parse Server via Back4App (parse_server_sdk_flutter)",
        "Service layer: AuthService, TaskService, CommentService (separates business logic)",
    ],
    notes="Point out where to configure Back4App keys (lib/services/back4app_config.dart).",
)

add_bulleted_slide(
    "Data Models",
    [
        "Task: title, description, status (New/Active/Closed), startDate, endDate, discussion",
        "TaskComment: threaded comments linked to Task (author, text, timestamps)",
        "ParseObjects used for persistence; services wrap Parse queries and saves",
    ],
    notes="Note the design choice: `discussion` field was used initially and converted to TaskComment entries when creating a task.",
)

add_bulleted_slide(
    "Key Screens / UX",
    [
        "Login/Register: autofill, password visibility toggle, validators",
        "Home: search with debounce, status filters, task cards with status pill and dates",
        "Add/Edit Task: expanded/constrained form, date pickers, validation (end >= start)",
        "Task Detail: threaded comments UI with edit/delete for owner",
    ],
    notes="Mention accessibility notes: central theme, inputDecorationTheme, improved color contrast.",
)

add_bulleted_slide(
    "Notable Features Implemented",
    [
        "Azure-DevOps-like task fields (status, startDate, endDate)",
        "Threaded comments and CommentService for CRUD",
        "Centralized theming and InputDecorationTheme for consistent inputs",
        "Keyboard safe layouts (AnimatedPadding + MediaQuery.viewInsets) to avoid clipping",
    ],
    notes="These are the main feature additions requested during the project.",
)

add_bulleted_slide(
    "Quality & Tooling",
    [
        "Static analysis (`flutter analyze`) and automated `dart fix` prefer_const fixes",
        "CI: GitHub Actions workflow added to run analyze/test/build",
        "Builds: web build validated (build/web produced) — Wasm warnings noted",
    ],
    notes="CI file: .github/workflows/flutter_ci.yml — runs analyze, tests, and web build.",
)

add_bulleted_slide(
    "Major Learnings",
    [
        "Centralize styling to avoid input inconsistencies and reduce duplication",
        "Small widget edits can introduce syntax balance issues; prefer incremental refactors",
        "Automated fixes (dart fix) are safe for prefer_const changes but review is advised",
        "Web builds may warn about packages using `dart:ffi` (not web-compatible) — review transitive deps",
    ],
    notes="These lessons are useful for future contributors and maintenance planning.",
)

add_bulleted_slide(
    "Run & Debug (How to)",
    [
        "Setup: Flutter SDK, appropriate device (Chrome) — see README.md",
        "Init: `flutter pub get`",
        "Run web: `flutter run -d chrome` — DevTools URL printed",
        "Build web: `flutter build web` (artifact at build/web)",
    ],
    notes="If you need Android/iOS builds, ensure platform SDKs are installed (Android SDK/Xcode).",
)

add_bulleted_slide(
    "Next Steps / Recommendations",
    [
        "Optional: full accessibility pass (semantics, contrast, keyboard navigation)",
        "Optional: dependency upgrades (run `flutter pub outdated` then incremental upgrades)",
        "Add integration tests to exercise auth/task/comment flows automatically",
    ],
    notes="I can help implement any of these next steps if you choose one.",
)

add_bulleted_slide(
    "Appendix: Key Files to Inspect",
    [
        "`lib/main.dart` — app entry and theme wiring",
        "`lib/theme/app_theme.dart` — centralized colors and InputDecorationTheme",
        "`lib/services/*` — AuthService, TaskService, CommentService",
        "`lib/screens/*` and `lib/widgets/*` — UI implementation",
    ],
    notes="Point developers to these files for quick onboarding.",
)

# Save
out_path = 'project_overview.pptx'
prs.save(out_path)
print('Wrote', out_path)
